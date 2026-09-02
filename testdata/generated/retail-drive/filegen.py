"""Realistic file-content generators for the OtterWorks enterprise drive.

Each function returns ``(bytes, mime_type)`` for a given logical file so the
uploaded objects are real, openable files of the correct type rather than empty
placeholders. Non-financial variety is derived deterministically from the file
name + a seeded ``random.Random`` so re-runs are reproducible; every financial
figure (prices, costs, margins, revenue) comes from the shared product catalog
(``catalog.py``, backed by the OTD-15 ``testdata/market-series/`` contract) so
all artifacts and the margins analytics dashboard show consistent numbers.

Heavy office formats (xlsx/docx/pptx/pdf/png/jpg) use optional third-party
libraries. If a library is missing the generator degrades gracefully to a
plain-text stand-in with the correct extension so the drive still populates.
"""
from __future__ import annotations

import io
import json
import math
import random
import struct
import wave
from datetime import date, datetime, timedelta
from decimal import Decimal

import catalog

# ---- optional heavy deps (degrade gracefully) -------------------------------
try:
    from openpyxl import Workbook
except Exception:  # pragma: no cover
    Workbook = None
try:
    import docx
except Exception:  # pragma: no cover
    docx = None
try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception:  # pragma: no cover
    Presentation = None
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas as _pdfcanvas
    from reportlab.platypus import (
        Image as RLImage,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )
except Exception:  # pragma: no cover
    _pdfcanvas = None
try:
    from PIL import Image, ImageDraw
except Exception:  # pragma: no cover
    Image = None
try:
    # Object-oriented API only (no pyplot): builders run in a thread pool and
    # pyplot's global figure registry is not thread-safe.
    from matplotlib.backends.backend_agg import FigureCanvasAgg
    from matplotlib.figure import Figure
except Exception:  # pragma: no cover
    Figure = None

MIME = {
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
    "csv": "text/csv",
    "txt": "text/plain",
    "md": "text/markdown",
    "json": "application/json",
    "png": "image/png",
    "jpg": "image/jpeg",
    "svg": "image/svg+xml",
    "mp4": "video/mp4",
    "wav": "audio/wav",
}

COMPANY = "OtterWorks"
_REGIONS = ["Northeast", "Southeast", "Midwest", "Southwest", "West", "Pacific Northwest"]
_STORES = [f"Store #{1000 + i}" for i in range(60)]

# Figures are sampled inside the committed baseline window so re-runs are
# stable regardless of "today" (dates past the baseline use the seeded walk).
_BASE_START = date(2024, 8, 1)
_BASE_END = date(2026, 6, 30)


def _rng(name: str, seed: int) -> random.Random:
    return random.Random(f"{seed}:{name}")


def _pick_date(r: random.Random) -> date:
    span = (_BASE_END - _BASE_START).days
    return _BASE_START + timedelta(days=r.randint(0, span))


def _pick_sku(r: random.Random) -> catalog.Sku:
    return r.choice(catalog.skus())


def sku_row(sku: catalog.Sku, d: date) -> dict[str, str]:
    """Canonical per-SKU figure row — the one source for every artifact."""
    cogs = catalog.cogs_usd(sku, d)
    return {
        "date": d.isoformat(),
        "sku": sku.sku,
        "product": sku.name,
        "category": sku.category,
        "supplier": sku.supplier,
        "commodity": sku.commodity_series_code,
        "list_price_usd": str(sku.list_price_usd),
        "cogs_usd": str(cogs),
        "margin_pct": str(catalog.margin_pct(sku, d)),
    }


# ---- individual format builders --------------------------------------------
def _xlsx(name: str, r: random.Random) -> bytes:
    if Workbook is None:
        return _txt_fallback(name, r)
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    headers = ["Date", "Region", "Store", "SKU", "Product", "Supplier", "Units",
               "List Price USD", "COGS USD", "Revenue USD", "Margin %"]
    ws.append(headers)
    for _ in range(r.randint(40, 200)):
        d = _pick_date(r)
        sku = _pick_sku(r)
        row = sku_row(sku, d)
        units = r.randint(1, 800)
        revenue = (Decimal(row["list_price_usd"]) * units).quantize(Decimal("0.01"))
        ws.append([
            row["date"], r.choice(_REGIONS), r.choice(_STORES), row["sku"],
            row["product"], row["supplier"], units, row["list_price_usd"],
            row["cogs_usd"], str(revenue), row["margin_pct"],
        ])
    s2 = wb.create_sheet("Summary")
    s2.append(["Metric", "Value"])
    s2.append(["Total Rows", ws.max_row - 1])
    s2.append(["Catalog SKUs", len(catalog.skus())])
    s2.append(["Generated", datetime.utcnow().isoformat()])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _docx(name: str, r: random.Random) -> bytes:
    if docx is None:
        return _txt_fallback(name, r)
    doc = docx.Document()
    doc.add_heading(name, level=0)
    doc.add_paragraph(
        f"{COMPANY} — Confidential. This document is part of the enterprise "
        "reference drive used for demonstration purposes."
    )
    for _ in range(r.randint(4, 9)):
        doc.add_heading(r.choice([
            "Executive Summary", "Objectives", "Scope", "Timeline",
            "Risks & Mitigations", "Budget", "Next Steps", "Appendix",
        ]), level=1)
        for _ in range(r.randint(2, 4)):
            doc.add_paragraph(_lorem(r, r.randint(30, 70)))
    sku = _pick_sku(r)
    row = sku_row(sku, _pick_date(r))
    doc.add_heading("Reference Figures", level=1)
    doc.add_paragraph(
        f"{row['product']} ({row['sku']}, {row['supplier']}): list price "
        f"${row['list_price_usd']}, COGS ${row['cogs_usd']}, margin {row['margin_pct']}%."
    )
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx(name: str, r: random.Random) -> bytes:
    if Presentation is None:
        return _txt_fallback(name, r)
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = name
    title_slide.placeholders[1].text = f"{COMPANY} — Internal Deck"
    for _ in range(r.randint(3, 6)):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = r.choice([
            "Market Overview", "Quarterly Performance", "Category Strategy",
            "Store Rollout", "Customer Insights", "Roadmap", "Financials",
        ])
        body = s.placeholders[1].text_frame
        body.text = _lorem(r, 12)
        sku = _pick_sku(r)
        row = sku_row(sku, _pick_date(r))
        p = body.add_paragraph()
        p.text = (f"• {row['product']}: ${row['list_price_usd']} list, "
                  f"{row['margin_pct']}% margin ({row['supplier']})")
        for _ in range(r.randint(1, 3)):
            p = body.add_paragraph()
            p.text = "• " + _lorem(r, r.randint(6, 12))
    # embedded image slide (chart if matplotlib is present, product art otherwise)
    img = _chart_png(name, r) if Figure else _product_art_png(name, r)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Market Chart"
    s.shapes.add_picture(io.BytesIO(img), Inches(1), Inches(1.5), width=Inches(8))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pdf(name: str, r: random.Random) -> bytes:
    if _pdfcanvas is None:
        return _txt_fallback(name, r)
    buf = io.BytesIO()
    c = _pdfcanvas.Canvas(buf, pagesize=LETTER)
    width, height = LETTER
    for page in range(r.randint(1, 4)):
        c.setFont("Helvetica-Bold", 16)
        c.drawString(72, height - 72, name if page == 0 else f"{name} (cont.)")
        c.setFont("Helvetica", 10)
        y = height - 110
        for _ in range(r.randint(20, 34)):
            c.drawString(72, y, _lorem(r, r.randint(8, 16)))
            y -= 16
            if y < 72:
                break
        c.showPage()
    c.save()
    return buf.getvalue()


def _sku_table(rows: list[dict[str, str]], style_rows=True) -> "Table":
    data = [["SKU", "Product", "Supplier", "List USD", "COGS USD", "Margin %"]]
    data += [[x["sku"], x["product"], x["supplier"], x["list_price_usd"],
              x["cogs_usd"], x["margin_pct"]] for x in rows]
    t = Table(data, repeatRows=1)
    style = [
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1e5f74")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
    ]
    if style_rows:
        style.append(("ROWBACKGROUNDS", (0, 1), (-1, -1),
                      [colors.whitesmoke, colors.HexColor("#e8f4f8")]))
    t.setStyle(TableStyle(style))
    return t


def _rich_pdf(name: str, r: random.Random, doc_type: str) -> bytes:
    """Multi-section platypus PDF (contract / spec sheet / invoice) with a
    catalog figure table and an embedded product-art image."""
    if _pdfcanvas is None:
        return _txt_fallback(name, r)
    styles = getSampleStyleSheet()
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=LETTER, title=name)
    d = _pick_date(r)
    skus = r.sample(catalog.skus(), k=min(6, len(catalog.skus())))
    rows = [sku_row(s, d) for s in skus]
    sections = {
        "contract": ["Parties & Term", "Supply Commitments", "Pricing Schedule",
                     "Quality & Compliance", "Termination"],
        "spec_sheet": ["Product Overview", "Materials & Sourcing",
                       "Cost Breakdown", "Packaging", "Compliance"],
        "invoice": ["Bill To", "Line Items", "Payment Terms"],
    }[doc_type]
    story = [Paragraph(name, styles["Title"]),
             Paragraph(f"{COMPANY} — Confidential", styles["Italic"]),
             Spacer(1, 0.2 * inch)]
    for section in sections:
        story.append(Paragraph(section, styles["Heading2"]))
        story.append(Paragraph(_lorem(r, r.randint(30, 60)), styles["BodyText"]))
        story.append(Spacer(1, 0.1 * inch))
    story.append(Paragraph(f"Catalog figures as of {d.isoformat()}", styles["Heading2"]))
    story.append(_sku_table(rows))
    story.append(Spacer(1, 0.2 * inch))
    art = _product_art_png(rows[0]["product"], r)
    story.append(RLImage(io.BytesIO(art), width=3 * inch, height=2 * inch))
    if doc_type == "invoice":
        total = sum(Decimal(x["list_price_usd"]) for x in rows)
        story.append(Paragraph(f"Total due: ${total}", styles["Heading3"]))
    doc.build(story)
    return buf.getvalue()


def _csv(name: str, r: random.Random) -> bytes:
    rows = ["date,region,store,sku,product,supplier,units,list_price_usd,cogs_usd,revenue_usd,margin_pct"]
    for _ in range(r.randint(50, 400)):
        d = _pick_date(r)
        sku = _pick_sku(r)
        row = sku_row(sku, d)
        units = r.randint(1, 900)
        revenue = (Decimal(row["list_price_usd"]) * units).quantize(Decimal("0.01"))
        rows.append(
            f"{row['date']},{r.choice(_REGIONS)},{r.choice(_STORES)},{row['sku']},"
            f"\"{row['product']}\",\"{row['supplier']}\",{units},"
            f"{row['list_price_usd']},{row['cogs_usd']},{revenue},{row['margin_pct']}"
        )
    return ("\n".join(rows) + "\n").encode()


def _json(name: str, r: random.Random) -> bytes:
    records = []
    for _ in range(r.randint(10, 60)):
        row = sku_row(_pick_sku(r), _pick_date(r))
        row["id"] = r.randint(10000, 99999)
        row["region"] = r.choice(_REGIONS)
        row["store"] = r.choice(_STORES)
        row["units"] = r.randint(1, 500)
        records.append(row)
    obj = {
        "name": name,
        "company": COMPANY,
        "generatedAt": datetime.utcnow().isoformat() + "Z",
        "records": records,
    }
    return json.dumps(obj, indent=2).encode()


def _md(name: str, r: random.Random) -> bytes:
    lines = [f"# {name}", "", f"> {COMPANY} internal reference document.", ""]
    for _ in range(r.randint(3, 6)):
        lines.append("## " + r.choice([
            "Overview", "Details", "Process", "Owners", "SLAs", "Checklist",
        ]))
        lines.append("")
        for _ in range(r.randint(3, 6)):
            lines.append("- " + _lorem(r, r.randint(6, 14)))
        lines.append("")
    return ("\n".join(lines)).encode()


def _txt(name: str, r: random.Random) -> bytes:
    return _lorem(r, r.randint(80, 240)).encode()


def _txt_fallback(name: str, r: random.Random) -> bytes:
    return (f"{name}\n\n" + _lorem(r, 120)).encode()


def _chart_png(name: str, r: random.Random) -> bytes:
    """Margin-trend / commodity-price chart from the shared market series."""
    if Figure is None:
        return _png(name, r)
    lower = name.lower()
    fig = Figure(figsize=(8, 4.5), dpi=100)
    FigureCanvasAgg(fig)
    ax = fig.add_subplot(111)
    if "margin" in lower:
        skus = r.sample(catalog.skus(), k=3)
        dates = [_BASE_START + timedelta(days=i) for i in range(0, 699, 14)]
        for sku in skus:
            ax.plot(dates, [float(catalog.margin_pct(sku, d)) for d in dates],
                    label=f"{sku.sku} {sku.name}")
        ax.set_ylabel("Margin %")
    else:
        code = "DREWRY_WCI_USD_FEU" if "freight" in lower else "SALMON_NOK_KG"
        s = catalog.series()[code]
        dates = [_BASE_START + timedelta(days=i) for i in range(0, 699, 7)]
        ax.plot(dates, [float(catalog.price_on(code, d)) for d in dates],
                color="#1e5f74", label=s.name)
        ax.set_ylabel(s.unit)
    ax.set_title(f"{name} — {COMPANY}")
    ax.legend(fontsize=8)
    ax.grid(alpha=0.3)
    fig.autofmt_xdate()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    return buf.getvalue()


_PALETTE = [
    (30, 95, 116), (72, 139, 73), (222, 143, 44), (149, 82, 145),
    (192, 78, 62), (54, 111, 168),
]


def _product_art_png(name: str, r: random.Random) -> bytes:
    """Stylized per-SKU product art (Pillow): badge + wave motif + label."""
    if Image is None:
        return _tiny_png()
    w, h = 640, 480
    bg = r.choice(_PALETTE)
    img = Image.new("RGB", (w, h), bg)
    d = ImageDraw.Draw(img)
    # water waves
    accent = tuple(min(255, c + 60) for c in bg)
    for i in range(6):
        y = 300 + i * 30
        for x in range(0, w, 40):
            d.arc([x, y, x + 40, y + 24], 180, 360, fill=accent, width=3)
    # product badge
    d.ellipse([190, 60, 450, 320], fill=(245, 240, 228), outline=accent, width=6)
    d.ellipse([250, 120, 390, 260], fill=r.choice(_PALETTE))
    d.text((210, 340), name[:44], fill=(255, 255, 255))
    d.text((210, 360), f"{COMPANY} — products for otters", fill=(230, 230, 230))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _svg(name: str, r: random.Random) -> bytes:
    """Per-SKU vector logo sticker (templated SVG; previewable as image/svg+xml)."""
    c1 = "#%02x%02x%02x" % r.choice(_PALETTE)
    c2 = "#%02x%02x%02x" % r.choice(_PALETTE)
    label = name[:36].replace("&", "&amp;").replace("<", "&lt;")
    svg = f"""<svg xmlns="http://www.w3.org/2000/svg" width="480" height="360" viewBox="0 0 480 360">
  <rect width="480" height="360" rx="24" fill="{c1}"/>
  <circle cx="240" cy="150" r="90" fill="#f5f0e4"/>
  <circle cx="240" cy="150" r="56" fill="{c2}"/>
  <path d="M40 300 q30 -24 60 0 t60 0 t60 0 t60 0 t60 0 t60 0" stroke="#f5f0e4" stroke-width="6" fill="none"/>
  <text x="240" y="330" text-anchor="middle" font-family="sans-serif" font-size="20" fill="#f5f0e4">{label}</text>
  <text x="240" y="40" text-anchor="middle" font-family="sans-serif" font-size="16" fill="#f5f0e4">{COMPANY}</text>
</svg>
"""
    return svg.encode()


def _png(name: str, r: random.Random) -> bytes:
    if Image is None:
        return _tiny_png()
    w, h = 640, 360
    img = Image.new("RGB", (w, h), (r.randint(20, 60), r.randint(40, 90), r.randint(80, 160)))
    d = ImageDraw.Draw(img)
    for _ in range(r.randint(8, 20)):
        x0, y0 = r.randint(0, w), r.randint(0, h)
        x1, y1 = x0 + r.randint(20, 180), y0 + r.randint(20, 120)
        d.rectangle([x0, y0, x1, y1], fill=(r.randint(0, 255), r.randint(0, 255), r.randint(0, 255)))
    d.text((16, 16), name[:40], fill=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _jpg(name: str, r: random.Random) -> bytes:
    if Image is None:
        return _tiny_png()
    w, h = 800, 600
    img = Image.new("RGB", (w, h), (r.randint(60, 200), r.randint(60, 200), r.randint(60, 200)))
    d = ImageDraw.Draw(img)
    for _ in range(r.randint(10, 25)):
        cx, cy = r.randint(0, w), r.randint(0, h)
        rad = r.randint(10, 90)
        d.ellipse([cx - rad, cy - rad, cx + rad, cy + rad],
                  fill=(r.randint(0, 255), r.randint(0, 255), r.randint(0, 255)))
    d.text((16, 16), name[:40], fill=(0, 0, 0))
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=82)
    return buf.getvalue()


def _wav(name: str, r: random.Random) -> bytes:
    """Build a deterministic, short mono tone/chirp for browser playback."""
    sample_rate = 8_000
    duration_seconds = 2
    frame_count = sample_rate * duration_seconds
    start_frequency = r.randint(240, 520)
    frequency_delta = r.randint(80, 260)
    amplitude = r.randint(8_000, 14_000)
    frames = bytearray()
    for index in range(frame_count):
        progress = index / frame_count
        frequency = start_frequency + frequency_delta * progress
        sample = int(amplitude * math.sin(2 * math.pi * frequency * index / sample_rate))
        frames.extend(struct.pack("<h", sample))

    buf = io.BytesIO()
    with wave.open(buf, "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)
        wav_file.setframerate(sample_rate)
        wav_file.writeframes(frames)
    return buf.getvalue()


def _tiny_png() -> bytes:
    return bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
        "0000000d49444154789c6360000002000100" "05fe02fea7d0e3450000000049454e44ae426082"
    )


_LOREM = (
    "revenue margin inventory assortment planogram markdown replenishment "
    "supplier logistics fulfillment omnichannel loyalty conversion basket "
    "shrinkage forecast promotion category compliance staffing footfall "
    "clearance seasonal warehouse distribution merchandising procurement audit "
    "salmon kelp tidepool river raft grooming whiskers pelt burrow estuary"
).split()


def _lorem(r: random.Random, n: int) -> str:
    words = [r.choice(_LOREM) for _ in range(n)]
    words[0] = words[0].capitalize()
    return " ".join(words) + "."


_BUILDERS = {
    "xlsx": _xlsx, "docx": _docx, "pptx": _pptx, "pdf": _pdf, "csv": _csv,
    "json": _json, "md": _md, "txt": _txt, "png": _png, "jpg": _jpg, "svg": _svg,
    "wav": _wav,
}

# kind -> builder variants (spec "kind" field in taxonomy.py).
_KIND_BUILDERS = {
    "chart": _chart_png,
    "product_art": _product_art_png,
    "contract": lambda n, r: _rich_pdf(n, r, "contract"),
    "spec_sheet": lambda n, r: _rich_pdf(n, r, "spec_sheet"),
    "invoice": lambda n, r: _rich_pdf(n, r, "invoice"),
}


def build(ext: str, name: str, seed: int, kind: str | None = None) -> tuple[bytes, str]:
    """Return (bytes, mime_type) for a file of type ``ext`` named ``name``."""
    ext = ext.lower()
    r = _rng(name, seed)
    builder = _KIND_BUILDERS.get(kind) or _BUILDERS.get(ext, _txt)
    data = builder(name, r)
    return data, MIME.get(ext, "application/octet-stream")
